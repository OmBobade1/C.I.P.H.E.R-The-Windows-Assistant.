import os
import speech_recognition as sr
import pyttsx3
from pptx import Presentation

# Initialize text-to-speech engine
engine = pyttsx3.init()

def speak(text):
    """Speak the given text using text-to-speech engine."""
    engine.say(text)
    engine.runAndWait()

def listen_for_command():
    """Listen for a voice command and return it as text."""
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print("Listening for command...")
        recognizer.adjust_for_ambient_noise(source)
        audio = recognizer.listen(source)
    
    try:
        command = recognizer.recognize_google(audio).lower()
        print(f"User said: {command}")
        return command
    except sr.UnknownValueError:
        print("Sorry, I didn't catch that.")
        speak("Sorry, I didn't catch that. Please try again.")
        return None
    except sr.RequestError:
        print("Sorry, there was an issue with the speech service.")
        speak("Sorry, there was an issue with the speech service.")
        return None

def create_presentation(command):
    """Create a new PowerPoint presentation with a given name."""
    pres_name = command.replace("create presentation", "").strip()
    
    if not pres_name:
        speak("Invalid name. Please say 'Create presentation' followed by a name.")
        return None, None

    pres_path = os.path.join(os.path.expanduser("~"), "Documents", f"{pres_name}.pptx")
    prs = Presentation()
    prs.save(pres_path)

    speak(f"Presentation '{pres_name}' has been created in Documents folder.")
    print(f"Presentation saved at {pres_path}")

    return prs, pres_path  #  Return both the presentation object and its path

def add_slide(prs):
    """Add a new slide with title and content."""
    if prs is None:
        speak("No active presentation. Please create one first.")
        return
    
    speak("What should be the title of the slide?")
    slide_title = listen_for_command()

    if not slide_title:
        speak("No title detected. Using 'Untitled Slide'.")
        slide_title = "Untitled Slide"

    speak("What content should I add to the slide?")
    slide_content = listen_for_command()

    if not slide_content:
        speak("No content detected. Using 'No content'.")
        slide_content = "No content"

    slide_layout = prs.slide_layouts[1]  # Title and Content Layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = slide_title
    content.text = slide_content

    speak(f"Slide titled '{slide_title}' has been added.")
    print(f"Slide '{slide_title}' added with content: {slide_content}")

def save_presentation(prs, pres_path):
    """Save the PowerPoint presentation."""
    if prs is None:
        speak("No active presentation to save.")
        return
    
    prs.save(pres_path)
    speak("Presentation has been saved.")
    print(f"Presentation saved at {pres_path}")

#  This function will only run when chipher enters PowerPoint mode.
def start_ppt_manager(chipher_instance):
    """Start the PowerPoint manager without interfering with chipher."""
    speak("PowerPoint manager activated. Say 'Create presentation [name]' to begin.")

    chipher_instance.presentation = None
    chipher_instance.pres_path = None

    while True:
        command = listen_for_command()
        if command:
            if "create presentation" in command:
                chipher_instance.presentation, chipher_instance.pres_path = create_presentation(command)
            elif "add slide" in command:
                if chipher_instance.presentation:
                    add_slide(chipher_instance.presentation)
                else:
                    speak("No active presentation. Say 'Create presentation' first.")
            elif "save presentation" in command:
                if chipher_instance.presentation:
                    save_presentation(chipher_instance.presentation, chipher_instance.pres_path)
                else:
                    speak("No active presentation to save.")
            elif "exit" in command:
                speak("Exiting PowerPoint manager.")
                break  #  This ensures it only exits PowerPoint mode, not chipher
